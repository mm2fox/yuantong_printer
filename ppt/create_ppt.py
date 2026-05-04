import asyncio
import os
import requests
import json
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCREENSHOT_DIR = "e:/Project/Print_tool/temple-management/ppt/screenshots"
PPT_OUTPUT = "e:/Project/Print_tool/temple-management/ppt/缘通寺院信息管理系统使用手册.pptx"
BASE_URL = "http://localhost:3001"
API_URL = "http://localhost:8000"

os.makedirs(SCREENSHOT_DIR, exist_ok=True)

def login_via_api():
    print("通过API登录...")
    try:
        response = requests.post(f"{API_URL}/api/auth/login", json={
            "username": "admin",
            "password": "admin123",
            "serverType": "internal"
        })
        data = response.json()
        if data.get("access_token"):
            print(f"登录成功! Token: {data['access_token'][:20]}...")
            return data["access_token"], data.get("user", {})
        else:
            print(f"登录失败: {data}")
            return None, None
    except Exception as e:
        print(f"API登录错误: {e}")
        return None, None

async def take_screenshots(token, user_info):
    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True, args=['--disable-web-security'])
        
        user_info_json = json.dumps(user_info)
        
        screenshots = {}
        
        try:
            print("正在截取登录页面...")
            login_context = await browser.new_context(viewport={'width': 1920, 'height': 1080})
            login_page = await login_context.new_page()
            await login_page.goto(f"{BASE_URL}/login", wait_until='load')
            await login_page.wait_for_timeout(2000)
            await login_page.screenshot(path=f"{SCREENSHOT_DIR}/01_login.png", full_page=False)
            screenshots['login'] = f"{SCREENSHOT_DIR}/01_login.png"
            await login_context.close()
            
            print("创建已登录的浏览器上下文...")
            context = await browser.new_context(viewport={'width': 1920, 'height': 1080})
            
            await context.add_init_script(f'''
                Object.defineProperty(window, 'localStorage', {{
                    value: {{
                        getItem: function(key) {{
                            const data = {{
                                'token': '{token}',
                                'userInfo': '{user_info_json}'
                            }};
                            return data[key] || null;
                        }},
                        setItem: function(key, value) {{
                            this._data[key] = value;
                        }},
                        _data: {{
                            'token': '{token}',
                            'userInfo': '{user_info_json}'
                        }},
                        removeItem: function(key) {{}},
                        clear: function() {{}}
                    }},
                    writable: false,
                    configurable: false
                }});
            ''')
            
            page = await context.new_page()
            
            print("正在截取法会记录查询页面...")
            await page.goto(f"{BASE_URL}/query/fahui", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/02_fahui_query.png", full_page=False)
            screenshots['fahui_query'] = f"{SCREENSHOT_DIR}/02_fahui_query.png"
            
            print("正在截取施主查询页面...")
            await page.goto(f"{BASE_URL}/query/shizhu", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/03_shizhu_query.png", full_page=False)
            screenshots['shizhu_query'] = f"{SCREENSHOT_DIR}/03_shizhu_query.png"
            
            print("正在截取法会登记页面...")
            await page.goto(f"{BASE_URL}/query/register", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/04_fahui_register.png", full_page=False)
            screenshots['fahui_register'] = f"{SCREENSHOT_DIR}/04_fahui_register.png"
            
            print("正在截取新增法会登记弹窗...")
            await page.evaluate('''() => {
                const btn = document.querySelector('.card-header button.el-button--primary');
                if (btn && btn.textContent.includes('新增登记')) {
                    btn.click();
                }
            }''')
            await page.wait_for_timeout(8000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/04a_fahui_register_dialog.png", full_page=False)
            screenshots['fahui_register_dialog'] = f"{SCREENSHOT_DIR}/04a_fahui_register_dialog.png"
            try:
                await page.click('.el-dialog__headerbtn')
            except:
                pass
            await page.wait_for_timeout(500)
            
            print("正在截取施主管理页面...")
            await page.goto(f"{BASE_URL}/shizhu", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/05_shizhu_manage.png", full_page=False)
            screenshots['shizhu_manage'] = f"{SCREENSHOT_DIR}/05_shizhu_manage.png"
            
            print("正在截取新增施主弹窗...")
            await page.evaluate('''() => {
                const btn = document.querySelector('.card-header button.el-button--primary');
                if (btn && btn.textContent.includes('新增施主')) {
                    btn.click();
                }
            }''')
            await page.wait_for_timeout(8000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/05a_shizhu_dialog.png", full_page=False)
            screenshots['shizhu_dialog'] = f"{SCREENSHOT_DIR}/05a_shizhu_dialog.png"
            try:
                await page.click('.el-dialog__headerbtn')
            except:
                pass
            await page.wait_for_timeout(500)
            
            print("正在截取法会管理页面...")
            await page.goto(f"{BASE_URL}/fahui", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/06_fahui_manage.png", full_page=False)
            screenshots['fahui_manage'] = f"{SCREENSHOT_DIR}/06_fahui_manage.png"
            
            print("正在截取新增法会弹窗...")
            await page.evaluate('''() => {
                const btn = document.querySelector('.card-header button.el-button--primary');
                if (btn && btn.textContent.includes('新增法会')) {
                    btn.click();
                }
            }''')
            await page.wait_for_timeout(8000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/06a_fahui_dialog.png", full_page=False)
            screenshots['fahui_dialog'] = f"{SCREENSHOT_DIR}/06a_fahui_dialog.png"
            try:
                await page.click('.el-dialog__headerbtn')
            except:
                pass
            await page.wait_for_timeout(500)
            
            print("正在截取打印管理页面...")
            await page.goto(f"{BASE_URL}/print", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/07_print_manage.png", full_page=False)
            screenshots['print_manage'] = f"{SCREENSHOT_DIR}/07_print_manage.png"
            
            print("正在截取打印模板页面...")
            await page.goto(f"{BASE_URL}/print/templates", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/08_print_template.png", full_page=False)
            screenshots['print_template'] = f"{SCREENSHOT_DIR}/08_print_template.png"
            
            print("正在截取新增打印模板弹窗...")
            await page.evaluate('''() => {
                const btn = document.querySelector('.card-header button.el-button--primary');
                if (btn && btn.textContent.includes('新增模板')) {
                    btn.click();
                }
            }''')
            await page.wait_for_timeout(8000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/08a_template_dialog.png", full_page=False)
            screenshots['template_dialog'] = f"{SCREENSHOT_DIR}/08a_template_dialog.png"
            
            print("正在截取打印模板编辑界面（滚动后）...")
            await page.evaluate('document.querySelector(".el-dialog__body").scrollTop = 300')
            await page.wait_for_timeout(500)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/08b_template_dialog2.png", full_page=False)
            screenshots['template_dialog2'] = f"{SCREENSHOT_DIR}/08b_template_dialog2.png"
            await page.click('.el-dialog__headerbtn')
            await page.wait_for_timeout(500)
            
            print("正在截取用户数据页面...")
            await page.goto(f"{BASE_URL}/system/user-data", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/09_user_data.png", full_page=False)
            screenshots['user_data'] = f"{SCREENSHOT_DIR}/09_user_data.png"
            
            print("正在截取用户管理页面...")
            await page.goto(f"{BASE_URL}/system/users", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/10_user_manage.png", full_page=False)
            screenshots['user_manage'] = f"{SCREENSHOT_DIR}/10_user_manage.png"
            
            print("正在截取新增用户弹窗...")
            await page.evaluate('''() => {
                const btn = document.querySelector('.card-header button.el-button--primary');
                if (btn && btn.textContent.includes('新增用户')) {
                    btn.click();
                }
            }''')
            await page.wait_for_timeout(8000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/10a_user_dialog.png", full_page=False)
            screenshots['user_dialog'] = f"{SCREENSHOT_DIR}/10a_user_dialog.png"
            try:
                await page.click('.el-dialog__headerbtn')
            except:
                pass
            await page.wait_for_timeout(500)
            
            print("正在截取寺庙管理页面...")
            await page.goto(f"{BASE_URL}/system/temples", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/11_temple_manage.png", full_page=False)
            screenshots['temple_manage'] = f"{SCREENSHOT_DIR}/11_temple_manage.png"
            
            print("正在截取系统日志页面...")
            await page.goto(f"{BASE_URL}/system/logs", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/12_system_log.png", full_page=False)
            screenshots['system_log'] = f"{SCREENSHOT_DIR}/12_system_log.png"
            
            print("正在截取数据库管理页面...")
            await page.goto(f"{BASE_URL}/system/database", wait_until='load')
            await page.wait_for_timeout(3000)
            await page.screenshot(path=f"{SCREENSHOT_DIR}/13_database.png", full_page=False)
            screenshots['database'] = f"{SCREENSHOT_DIR}/13_database.png"
            
        except Exception as e:
            print(f"截图过程中出错: {e}")
            import traceback
            traceback.print_exc()
        finally:
            await browser.close()
        
        return screenshots

def create_ppt(screenshots):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(102, 126, 234)
        shape.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items, screenshot_path=None):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(102, 126, 234)
        line.line.fill.background()
        
        if screenshot_path and os.path.exists(screenshot_path):
            slide.shapes.add_picture(screenshot_path, Inches(0.5), Inches(1.3), width=Inches(8))
            content_box = slide.shapes.add_textbox(Inches(8.7), Inches(1.3), Inches(4.3), Inches(5.5))
        else:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(12)
        
        return slide
    
    print("正在创建PPT...")
    
    add_title_slide("缘通寺院信息管理系统", "使用手册")
    
    add_content_slide("系统概述", [
        "缘通寺院信息管理系统是一款专为寺院设计的综合管理软件",
        "采用前后端分离架构，前端使用Vue3 + Element Plus",
        "后端使用FastAPI + SQLite，轻量级部署",
        "支持多用户、多寺庙管理",
        "提供法会登记、施主管理、打印模板等核心功能",
        "界面简洁美观，操作便捷高效"
    ])
    
    add_content_slide("系统架构", [
        "前端技术栈：Vue 3 + Vite + Element Plus + Pinia",
        "后端技术栈：FastAPI + SQLAlchemy + SQLite",
        "打印功能：支持自定义模板，可连接扫描仪",
        "安全机制：JWT认证 + 权限管理",
        "部署方式：支持独立部署，可打包为单机应用"
    ])
    
    add_content_slide("一、登录系统", [
        "选择用户：从下拉列表选择用户名",
        "输入密码：输入对应用户的密码",
        "默认管理员账号：admin / admin123"
    ], screenshots.get('login'))
    
    add_content_slide("二、法会记录查询", [
        "多条件搜索：法会名称、施主姓名、编号等",
        "日期范围筛选：支持开始和结束日期",
        "牌位类型筛选：大牌、中牌、小牌",
        "类型筛选：延生/往生",
        "打印状态筛选：已打印/未打印",
        "支持批量标记打印状态",
        "支持导出CSV文件"
    ], screenshots.get('fahui_query'))
    
    add_content_slide("三、施主查询", [
        "按施主姓名或编号搜索",
        "查看施主的详细法会记录",
        "支持查看历史登记信息",
        "快速定位施主相关数据"
    ], screenshots.get('shizhu_query'))
    
    add_content_slide("四、法会登记", [
        "新增法会登记记录",
        "选择法会：从已有法会中选择或新增",
        "选择施主：搜索选择或新增施主",
        "填写姓名信息：最多5个姓名",
        "选择牌位类型：大牌、中牌、小牌",
        "选择类型：延生或往生",
        "填写金额和登记日期",
        "支持编辑和删除记录"
    ], screenshots.get('fahui_register'))
    
    add_content_slide("四、法会登记 - 新增登记弹窗", [
        "法会选择：下拉选择已有法会",
        "施主选择：搜索或新增施主",
        "姓名填写：最多填写5个姓名",
        "牌位类型：大牌、中牌、小牌",
        "类型选择：延生或往生",
        "金额输入：功德金金额",
        "登记日期：默认当天，可修改"
    ], screenshots.get('fahui_register_dialog'))
    
    add_content_slide("五、施主管理", [
        "新增、编辑、删除施主信息",
        "施主编号自动生成",
        "记录联系方式和地址",
        "功德主标记功能",
        "佛光接引、阳上、佛光注照信息管理",
        "查看施主的法会记录历史"
    ], screenshots.get('shizhu_manage'))
    
    add_content_slide("五、施主管理 - 新增施主弹窗", [
        "施主姓名：必填项",
        "联系电话：可选填写",
        "地址：可选填写",
        "功德主：是否标记为功德主",
        "佛光接引、阳上、佛光注照：往生牌位信息",
        "支持双击表格行快速编辑姓名"
    ], screenshots.get('shizhu_dialog'))
    
    add_content_slide("六、法会管理", [
        "管理所有法会信息",
        "设置法会名称、日期",
        "设置功德金标准（大/中/小牌）",
        "查看法会登记统计"
    ], screenshots.get('fahui_manage'))
    
    add_content_slide("六、法会管理 - 新增法会弹窗", [
        "法会名称：如'华严法会'",
        "开始日期：法会开始时间",
        "结束日期：法会结束时间",
        "功德金标准：设置大/中/小牌的功德金",
        "备注：可选填写"
    ], screenshots.get('fahui_dialog'))
    
    add_content_slide("七、打印管理", [
        "延生法会打印：打印延生牌位",
        "往生法会打印：打印往生牌位",
        "支持选择打印模板",
        "批量打印功能",
        "打印预览功能"
    ], screenshots.get('print_manage'))
    
    add_content_slide("八、打印模板管理（管理员）", [
        "创建和管理打印模板",
        "支持扫描仪扫描底图",
        "自定义页面尺寸",
        "自定义字体和字号",
        "姓名区域位置调整",
        "阳上区域设置（往生牌位）",
        "预览和打印测试"
    ], screenshots.get('print_template'))
    
    add_content_slide("八、打印模板 - 新增/编辑模板弹窗", [
        "模板名称：自定义模板名称",
        "模板类型：延生牌位/往生牌位/佛事牌子",
        "牌位类型：大牌/中牌/小牌",
        "是否启用：控制模板是否可用",
        "扫描底图：支持扫描仪或上传图片"
    ], screenshots.get('template_dialog'))
    
    add_content_slide("八、打印模板 - 页面与姓名设置", [
        "页面设置：自定义宽度和高度（毫米）",
        "字体选择：华文行楷、宋体、黑体、楷体",
        "姓名字号：调整姓名字体大小",
        "姓名间距：多个姓名之间的间距",
        "区域上边距：姓名区域距离顶部的位置",
        "底图透明度：调整底图显示透明度"
    ], screenshots.get('template_dialog2'))
    
    add_content_slide("九、用户数据", [
        "查看个人操作记录",
        "查看个人统计数据",
        "管理个人相关信息"
    ], screenshots.get('user_data'))
    
    add_content_slide("十、用户管理（管理员）", [
        "新增、编辑、删除用户",
        "设置用户角色：管理员/普通用户",
        "分配用户所属寺庙",
        "重置用户密码",
        "设置用户权限",
        "启用/禁用用户"
    ], screenshots.get('user_manage'))
    
    add_content_slide("十、用户管理 - 新增用户弹窗", [
        "用户名：登录账号",
        "密码：初始密码",
        "姓名：用户真实姓名",
        "角色：管理员或普通用户",
        "所属寺庙：分配用户管理的寺庙",
        "是否启用：控制用户是否可登录"
    ], screenshots.get('user_dialog'))
    
    add_content_slide("十一、寺庙管理（管理员）", [
        "新增、编辑、删除寺庙",
        "设置寺庙名称、地址、联系电话",
        "查看各寺庙用户数量",
        "有用户的寺庙无法删除"
    ], screenshots.get('temple_manage'))
    
    add_content_slide("十二、系统日志（管理员）", [
        "查看系统操作日志",
        "记录用户登录、操作等行为",
        "支持日志筛选和查询",
        "便于系统审计和问题追踪"
    ], screenshots.get('system_log'))
    
    add_content_slide("十三、数据库管理（管理员）", [
        "数据库备份功能",
        "数据库恢复功能",
        "查看数据库状态",
        "数据导出功能"
    ], screenshots.get('database'))
    
    add_content_slide("常见问题解答", [
        "Q: 忘记密码怎么办？",
        "A: 请联系管理员重置密码",
        "",
        "Q: 如何添加新法会？",
        "A: 在法会登记页面点击法会名称旁的'新增'按钮",
        "",
        "Q: 如何自定义打印模板？",
        "A: 进入打印模板管理，可以扫描已有牌位作为底图参考",
        "",
        "Q: 数据存储在哪里？",
        "A: 数据存储在SQLite数据库中，位于backend目录下"
    ])
    
    add_content_slide("技术支持", [
        "系统版本：v1.0.0",
        "开发框架：Vue 3 + FastAPI",
        "数据库：SQLite",
        "",
        "如遇问题，请联系系统管理员",
        "或查看项目文档获取更多帮助"
    ])
    
    add_title_slide("感谢使用", "缘通寺院信息管理系统")
    
    prs.save(PPT_OUTPUT)
    print(f"PPT已保存到: {PPT_OUTPUT}")

async def main():
    token, user_info = login_via_api()
    if not token:
        print("登录失败，无法继续截图")
        return
    
    print("开始截取系统截图...")
    screenshots = await take_screenshots(token, user_info)
    
    print(f"共截取 {len(screenshots)} 张截图")
    
    create_ppt(screenshots)
    print("完成!")

if __name__ == "__main__":
    asyncio.run(main())
